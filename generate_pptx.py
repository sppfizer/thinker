"""
PRM – Physical Routing Model
PowerPoint generator — run with: py generate_pptx.py
Outputs: PRM-Design.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2   = RGBColor(0xFF, 0xA5, 0x00)   # amber
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCA, 0xE9, 0xFF)
SUBTLE    = RGBColor(0x55, 0x77, 0x99)
GREEN     = RGBColor(0x06, 0xD6, 0xA0)
RED       = RGBColor(0xEF, 0x47, 0x6F)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # totally blank


# ── helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, l, t, w, h,
          size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def rect(slide, l, t, w, h, fill=ACCENT, alpha=None, line=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def divider(slide, y, color=ACCENT):
    ln = slide.shapes.add_connector(1, Inches(0.5), Inches(y), Inches(12.83), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(1)

def pill(slide, text, l, t, w=1.6, h=0.38, bg_col=ACCENT, txt_col=BG, size=13, bold=True):
    r = rect(slide, l, t, w, h, fill=bg_col)
    txbox(slide, text, l+0.05, t+0.03, w-0.1, h-0.06,
          size=size, bold=bold, color=txt_col, align=PP_ALIGN.CENTER)

def bullet_block(slide, items, l, t, w, h, size=16, color=WHITE, gap=0.33):
    y = t
    for item in items:
        prefix = "▸  " if not item.startswith("  ") else item[:3]
        body   = item.lstrip()
        txbox(slide, ("▸  " if not item.startswith("  ") else "    ◦  ") + body,
              l, y, w, gap, size=size, color=color)
        y += gap

def section_header(slide, number, title):
    rect(slide, 0, 0, 13.33, 0.6, fill=ACCENT)
    txbox(slide, f"{number}  {title}", 0.3, 0.08, 12, 0.5,
          size=22, bold=True, color=BG, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
rect(s, 0, 2.6, 13.33, 2.4, fill=RGBColor(0x05, 0x2B, 0x45))
txbox(s, "🧠  Physical Routing Model", 0.6, 2.7, 12, 1.0,
      size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txbox(s, "PRM  —  A non-neural-network language model", 0.6, 3.65, 12, 0.55,
      size=22, color=LIGHT, align=PP_ALIGN.CENTER)
txbox(s, "Forward-only training · No backpropagation · Physics-based routing",
      0.6, 4.2, 12, 0.5, size=16, color=SUBTLE, align=PP_ALIGN.CENTER, italic=True)
divider(s, 5.1)
txbox(s, "sppfizer  +  GitHub Copilot  ·  2026-07-08", 0.6, 5.2, 12, 0.4,
      size=13, color=SUBTLE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Problem
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "01", "Why Not Neural Networks?")

txbox(s, "Goal: a language model that trains fast, forward-only, with no backpropagation",
      0.5, 0.8, 12.3, 0.5, size=18, color=LIGHT, italic=True)

headers = ["Classical NN", "What We Reject", "Why"]
col_w = [2.8, 3.5, 5.5]
col_x = [0.4, 3.3, 6.9]
for i,(h,x,w) in enumerate(zip(headers, col_x, col_w)):
    rect(s, x, 1.45, w-0.1, 0.4, fill=ACCENT)
    txbox(s, h, x+0.1, 1.48, w-0.2, 0.35, size=14, bold=True, color=BG)

rows = [
    ("Backpropagation",    "Two-phase training",       "Forward pass stores all activations; backward pass expensive"),
    ("Gradient descent",   "Global loss computation",  "Must finish full pass before any weight update"),
    ("Learned embeddings", "Separate pre-training",    "Extra cost just to encode token identity"),
    ("Attention matrices", "O(n²) dot-product",        "Quadratic cost for every token-pair relationship"),
]
for i,(a,b,c) in enumerate(rows):
    y = 1.95 + i*0.55
    fc = RGBColor(0x12,0x2A,0x40) if i%2==0 else RGBColor(0x0A,0x20,0x33)
    for x,w,txt in zip(col_x, col_w, [a,b,c]):
        rect(s, x, y, w-0.1, 0.48, fill=fc)
        txbox(s, txt, x+0.1, y+0.06, w-0.2, 0.38, size=13, color=WHITE if i%2==0 else LIGHT)

txbox(s, "✦  Instead: one forward pass, the correct answer acts as a magnet guiding the model in real time",
      0.4, 4.65, 12.3, 0.5, size=16, bold=True, color=ACCENT2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — The Metaphor
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "02", "The Nail Grid Metaphor")

# Left: concept table
rect(s, 0.3, 0.75, 6.0, 0.38, fill=ACCENT)
for lbl, x in [("Physical Object", 0.35), ("Represents", 2.45), ("Role in PRM", 4.55)]:
    txbox(s, lbl, x, 0.78, 2.0, 0.32, size=13, bold=True, color=BG)

rows_m = [
    ("Ball",           "Token",                "Unit of meaning in input"),
    ("Ball size/mass", "Token frequency",      "Pre-computed weight — not learned"),
    ("Nail tilt",      "Routing weight",       "Primary learned parameter"),
    ("Nail diameter",  "Routing bias",         "Resistance to being nudged"),
    ("Magnet",         "Correct output token", "Forward correction force during training"),
    ("Vocab slot width","Token frequency",     "Built-in language prior"),
]
for i,(a,b,c) in enumerate(rows_m):
    y = 1.18 + i*0.48
    fc = RGBColor(0x12,0x2A,0x40) if i%2==0 else RGBColor(0x0A,0x20,0x33)
    rect(s, 0.3, y, 6.0, 0.44, fill=fc)
    for x,w,txt in zip([0.35,2.45,4.55],[2.05,2.05,1.7],[a,b,c]):
        txbox(s, txt, x, y+0.06, w, 0.34, size=12, color=WHITE)

# Right: ASCII grid diagram
rect(s, 6.7, 0.75, 6.2, 5.8, fill=RGBColor(0x08,0x1E,0x30))
diagram = [
    "  INPUT TOKENS",
    "  ● ●● ○ ● ○",
    "  ↓ ↓↓ ↓ ↓ ↓",
    "  ╔═══════════╗",
    "  ║ ⊕ ⊕ ⊕ ⊕ ║ ← nails",
    "  ║  ⊕ ⊕ ⊕  ║",
    "  ║ ⊕ ⊕ ⊕ ⊕ ║",
    "  ║  ⊕ ⊕ ⊕  ║",
    "  ╚═══════════╝",
    "  [the][a][cat]",
    "   ▓▓▓  ▓  ▓▓",
    "  wide  sm  med",
    "",
    "  ⬛ = magnet (training)",
    "  ● = heavy ball (freq)",
    "  ○ = light ball (rare)",
]
for i, line in enumerate(diagram):
    txbox(s, line, 6.8, 0.82 + i*0.34, 5.8, 0.35,
          size=12, color=ACCENT if i in [0,13,14,15] else (ACCENT2 if i==9 else WHITE))

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Ball-to-Ball Interaction
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "03", "Token Interaction — Physics Between Balls")

for col, (title, col_fill, items) in enumerate([
    ("🌌  Gravitational Attraction  (long range)", ACCENT, [
        "All balls pull each other: F ∝ mass × proximity",
        "Heavier balls exert stronger pull",
        "Naturally clusters semantically related tokens",
        '"Paris" curves toward "France" without any',
        "   learned attention matrix",
        "",
        "Creates semantic neighbourhoods",
        "   purely from physics",
    ]),
    ("💥  Elastic Collision  (short range)", ACCENT2, [
        "When ball paths cross → they collide",
        "Larger ball deflects less (more inertia)",
        "Smaller ball deflects more",
        "Dominant tokens (subject, verb) steer",
        "   weaker tokens (prepositions, fillers)",
        "",
        "Creates natural token dominance",
        "   hierarchy — no learned weights needed",
    ]),
]):
    x = 0.4 + col * 6.5
    rect(s, x, 0.75, 6.1, 0.42, fill=col_fill)
    txbox(s, title, x+0.1, 0.78, 5.9, 0.36, size=14, bold=True, color=BG)
    for i,item in enumerate(items):
        txbox(s, ("▸  " if item and not item.startswith(" ") else "") + item,
              x+0.15, 1.25+i*0.41, 5.8, 0.38, size=13,
              color=WHITE if item else SUBTLE)

divider(s, 5.2)
txbox(s, "Prediction Slot Ball  →  zero-mass ball dropped at 'next token' position."
         "  Gets shaped entirely by context balls.  Where it lands = the predicted token.",
      0.4, 5.3, 12.4, 0.5, size=15, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)
pill(s, "Physical analog of the Transformer Query vector", 3.5, 5.88, 6.3, 0.42,
     bg_col=RGBColor(0x05,0x3A,0x5C), txt_col=ACCENT, size=13, bold=False)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Forward Training / Magnet
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "04", "Forward-Only Training — The Magnet")

rect(s, 0.3, 0.75, 5.8, 5.5, fill=RGBColor(0x08,0x1E,0x30))
txbox(s, "Classical Backpropagation", 0.35, 0.8, 5.6, 0.4, size=15, bold=True, color=RED)
bp = [
    "① Full forward pass",
    "② Store ALL intermediate activations",
    "③ Compute global loss at output",
    "④ Propagate gradients BACKWARD",
    "   through every layer",
    "⑤ Update weights",
    "",
    "→ Two phases",
    "→ Memory: O(depth × batch)",
    "→ No update until pass completes",
]
for i,l in enumerate(bp):
    txbox(s, l, 0.45, 1.28+i*0.42, 5.5, 0.4, size=13, color=LIGHT if "→" in l else WHITE)

rect(s, 6.6, 0.75, 6.2, 5.5, fill=RGBColor(0x05,0x2B,0x1A))
txbox(s, "PRM Forward Magnetic Training", 6.65, 0.8, 6.0, 0.4, size=15, bold=True, color=GREEN)
fm = [
    "① Drop balls + prediction slot ball",
    "② Magnet at target slot — force felt",
    "   through entire diamond depth",
    "③ Each nail updates AS ball passes:",
    "",
    "   tilt += mass × magnet_force",
    "           × (1 / diameter)",
    "",
    "→ ONE phase — update is live",
    "→ Memory: O(1) per nail",
    "→ Rare answer = stronger magnet",
    "   = larger update = learns harder",
]
for i,l in enumerate(fm):
    c = ACCENT2 if "tilt +=" in l or "×" in l else (GREEN if "→" in l else WHITE)
    txbox(s, l, 6.7, 1.28+i*0.38, 6.0, 0.36, size=13, color=c)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Diamond Grid
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "05", "Diamond Grid Shape — Think · Mix · Summarize · Forget")

# Left: diagram
rect(s, 0.3, 0.75, 5.5, 6.1, fill=RGBColor(0x08,0x1E,0x30))
diag = [
    "  ● ● ○ ● ●    ← input tokens",
    "  ┌──────────┐  ← narrow entry",
    "   \\        /",
    "    \\  ⊕⊕  /   WIDENING",
    "     \\⊕⊕⊕/    (divergent",
    "      ⊕⊕⊕      thinking)",
    "     /⊕⊕⊕\\",
    "    /  ⊕⊕  \\   NARROWING",
    "   /        \\  (summarizing)",
    "  └──────────┘  ← narrow output",
    "  [the][a][cat] ← vocab slots",
]
for i,line in enumerate(diag):
    c = ACCENT if i in [0,10] else (ACCENT2 if "WIDENING" in line or "NARROWING" in line else WHITE)
    txbox(s, line, 0.4, 0.85+i*0.49, 5.2, 0.47, size=12, color=c)

# Right: properties
props = [
    ("Open Borders", GREEN,
     ["Balls drifting to edge fall off",
      "= natural information pruning",
      "Replaces: dropout, regularisation,",
      "  attention masking — all emergent"]),
    ("W:N Ratio", ACCENT,
     ["Widening : Narrowing rows",
      "1:1 = balanced think/summarize",
      "3:1 = deep exploration (creative)",
      "1:3 = quick spread, long focus"]),
    ("Memory ∝ Depth", ACCENT2,
     ["More rows = more nails",
      "More nails = more learned routing",
      "= higher context capacity",
      "Open borders keep output clean"]),
]
for i,(title, col, items) in enumerate(props):
    y = 0.75 + i*2.0
    rect(s, 6.1, y, 6.8, 1.85, fill=RGBColor(0x08,0x1E,0x30))
    rect(s, 6.1, y, 6.8, 0.38, fill=col)
    txbox(s, title, 6.2, y+0.04, 6.6, 0.32, size=14, bold=True, color=BG)
    for j,item in enumerate(items):
        txbox(s, "▸  "+item, 6.2, y+0.45+j*0.34, 6.6, 0.32, size=13, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Visual Mechanics
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "06A", "Visual Mechanics — Balls, Nails, Magnet Paths")

visual_panels = [
    ("Ball size = token weight", ACCENT, 0.35, 0.85, 3.95, 5.95),
    ("Nail thickness = bias/stability", ACCENT2, 4.45, 0.85, 4.25, 5.95),
    ("Magnet follows the diamond", GREEN, 8.9, 0.85, 4.05, 5.95),
]
for title, col, x, y, w, h in visual_panels:
    rect(s, x, y, w, h, fill=RGBColor(0x08,0x1E,0x30))
    rect(s, x, y, w, 0.4, fill=col)
    txbox(s, title, x+0.1, y+0.05, w-0.2, 0.3, size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)

# Ball size panel
txbox(s, "BIG BALL", 0.55, 1.35, 1.6, 0.3, size=12, bold=True, color=ACCENT2)
txbox(s, "●", 0.55, 1.65, 1.6, 1.0, size=56, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
txbox(s, "heavy token\nhigh frequency\nstronger pull", 0.55, 2.75, 1.6, 0.9, size=12, color=LIGHT, align=PP_ALIGN.CENTER)
txbox(s, "SMALL BALL", 2.55, 1.35, 1.6, 0.3, size=12, bold=True, color=GREEN)
txbox(s, "○", 2.55, 1.65, 1.6, 1.0, size=56, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txbox(s, "rare token\nlow frequency\nlighter path", 2.55, 2.75, 1.6, 0.9, size=12, color=LIGHT, align=PP_ALIGN.CENTER)
rect(s, 0.55, 4.1, 3.3, 1.1, fill=RGBColor(0x05,0x28,0x45))
txbox(s, "Massive balls hit nails harder\nand influence routing more.", 0.65, 4.22, 3.1, 0.9, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Nail thickness panel
txbox(s, "THICK NAIL", 4.65, 1.35, 1.8, 0.3, size=12, bold=True, color=ACCENT2)
txbox(s, "┃", 4.85, 1.65, 1.4, 1.2, size=68, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
txbox(s, "high bias\nstable routing\nhard to move", 4.65, 2.85, 1.8, 1.0, size=12, color=LIGHT, align=PP_ALIGN.CENTER)
txbox(s, "THIN NAIL", 6.75, 1.35, 1.6, 0.3, size=12, bold=True, color=GREEN)
txbox(s, "|", 7.05, 1.65, 0.9, 1.2, size=68, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txbox(s, "low bias\nflexible routing\neasy to nudge", 6.75, 2.85, 1.6, 1.0, size=12, color=LIGHT, align=PP_ALIGN.CENTER)
rect(s, 4.65, 4.1, 3.8, 1.1, fill=RGBColor(0x05,0x28,0x45))
txbox(s, "Thicker nails resist change.\nThinner nails learn faster.", 4.75, 4.22, 3.6, 0.9, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Magnet path panel
txbox(s, "WIDEN", 9.1, 1.35, 0.8, 0.3, size=12, bold=True, color=ACCENT)
txbox(s, "↙  ↓  ↘", 9.1, 1.65, 2.8, 0.6, size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txbox(s, "MID", 9.1, 2.4, 0.8, 0.3, size=12, bold=True, color=ACCENT2)
txbox(s, "←  ←  ↓  →  →", 9.1, 2.7, 3.0, 0.6, size=22, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
txbox(s, "NARROW", 9.1, 3.45, 1.1, 0.3, size=12, bold=True, color=GREEN)
txbox(s, "↘  ↓  ↙", 9.1, 3.75, 2.8, 0.6, size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
rect(s, 9.2, 4.5, 2.8, 0.95, fill=RGBColor(0x05,0x28,0x45), line=GREEN)
txbox(s, "Magnet fans outward\nthen focuses inward\nto the target slot.", 9.3, 4.62, 2.6, 0.75, size=12, color=WHITE, align=PP_ALIGN.CENTER)

divider(s, 6.95)
txbox(s, "Visual rule:  big ball + thick nail + focused magnet path = stronger, more stable routing signal",
      0.4, 7.0, 12.4, 0.35, size=13, color=LIGHT, italic=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Role Specialists
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "07", "Role-Based Specialist Diamonds")

txbox(s, "Specialists are trained on cognitive ROLES, not knowledge domains",
      0.4, 0.75, 12.4, 0.45, size=17, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)

rect(s, 0.3, 1.3, 5.9, 0.38, fill=RED)
txbox(s, "❌  Category Specialist  (brittle)", 0.4, 1.33, 5.7, 0.32, size=13, bold=True, color=WHITE)
rect(s, 6.5, 1.3, 6.4, 0.38, fill=GREEN)
txbox(s, "✅  Role Specialist  (generalises)", 6.6, 1.33, 6.2, 0.32, size=13, bold=True, color=BG)

for i,(l,r) in enumerate([
    ("Trained on: science data", "Trained on: analytical content across ALL domains"),
    ("Fails on: new science not in training", "Generalises to: any task requiring analytical thinking"),
    ('Router asks: "what topic is this?"', 'Router asks: "what thinking style is needed?"'),
    ("Brittle at domain borders", "Robust — roles transfer across all domains"),
]):
    y = 1.75+i*0.5
    fc = RGBColor(0x20,0x0A,0x0A) if i%2==0 else RGBColor(0x18,0x06,0x06)
    rect(s, 0.3, y, 5.9, 0.46, fill=fc); txbox(s, l, 0.4, y+0.08, 5.7, 0.34, size=13, color=LIGHT)
    fc2 = RGBColor(0x05,0x20,0x12) if i%2==0 else RGBColor(0x03,0x16,0x0C)
    rect(s, 6.5, y, 6.4, 0.46, fill=fc2); txbox(s, r, 6.6, y+0.08, 6.2, 0.34, size=13, color=WHITE)

divider(s, 3.9)
roles = [
    ("Analyst",         "1:1","Thick","Logical, step-by-step deduction"),
    ("Generator",       "3:1","Thin", "Creative, broad exploration"),
    ("Synthesizer",     "1:3","Med",  "Compression, abstraction"),
    ("Precisionist",    "1:4","V.Thick","Exact recall, minimal inference"),
    ("Narrator",        "2:2","Mixed","Sequential explanation, teaching"),
    ("Conversationalist","2:1","Thin","Context-aware, tone-matching"),
]
cols = ["Role", "W:N", "Nail ⌀", "Cognitive Style"]
col_x2 = [0.4, 3.4, 5.0, 6.5]
col_w2 = [2.8, 1.4, 1.3, 6.3]
rect(s, 0.3, 4.05, 12.7, 0.38, fill=ACCENT)
for h,x,w in zip(cols, col_x2, col_w2):
    txbox(s, h, x, 4.08, w, 0.32, size=13, bold=True, color=BG)
for i,(role,wn,nd,style) in enumerate(roles):
    y = 4.5+i*0.45
    fc = RGBColor(0x0D,0x25,0x3A) if i%2==0 else RGBColor(0x08,0x1C,0x2C)
    rect(s, 0.3, y, 12.7, 0.42, fill=fc)
    for txt,x,w in zip([role,wn,nd,style], col_x2, col_w2):
        txbox(s, txt, x, y+0.06, w, 0.32, size=12, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Routing / Specialists
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "08", "Routing Strategy — Selecting the Right Specialist")

for col,(title,col_fill,items,note) in enumerate([
    ("Parallel  (≤ 3 specialists)", ACCENT, [
        "Run all N diamonds simultaneously",
        "Same input tokens → each specialist",
        "Winner = highest retained ball-mass ratio",
        "   (balls falling off = out-of-role signal)",
        "Cost: N × S",
        "Best when N is small (2–3)",
    ], "retained mass = confidence"),
    ("Pre-Selector Diamond  (> 3 specialists)", ACCENT2, [
        "Small routing diamond runs first",
        "Output slots = specialist IDs",
        "Slot width ∝ activation frequency",
        "Winner gets input re-dropped identically",
        "Cost: R + S   (R ≈ 0.3 × S typically)",
        "Best when N > (R/S) + 1",
    ], "soft routing: border landing → two specialists activate"),
]):
    x = 0.4 + col*6.5
    rect(s, x, 0.75, 6.1, 0.42, fill=col_fill)
    txbox(s, title, x+0.1, 0.78, 5.9, 0.36, size=14, bold=True, color=BG)
    for i,item in enumerate(items):
        txbox(s, ("▸  " if item and not item.startswith(" ") else "")+item,
              x+0.15, 1.25+i*0.52, 5.8, 0.48, size=13, color=WHITE)
    rect(s, x, 4.4, 6.1, 0.38, fill=RGBColor(0x05,0x28,0x45))
    txbox(s, "ℹ  "+note, x+0.1, 4.43, 5.9, 0.32, size=12, color=LIGHT, italic=True)

divider(s, 5.0)
txbox(s, "EM Bootstrap (no labels):  train all specialists blindly → measure winner per sample → auto-label → train router → iterate",
      0.4, 5.1, 12.4, 0.48, size=13, color=SUBTLE, italic=True)
txbox(s, "Manual role labeling recommended for production — role = cognitive style, not topic domain",
      0.4, 5.6, 12.4, 0.45, size=14, bold=True, color=ACCENT2)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Math Model
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "09", "Mathematical Formulation")

rect(s, 0.3, 0.75, 12.7, 5.85, fill=RGBColor(0x05,0x18,0x2A))
math_lines = [
    ("State",           "Ball i at row r:   (x_i,  v_i,  m_i)   — position, velocity, mass"),
    ("Nail",            "Nail (r,c):        tilt ∈ [-1,1]   ·   diameter ∈ (0,1]"),
    ("Deflection",      "Δx_i  =  tilt[r,c]  ·  α  /  m_i          (inertia: heavier → less deflect)"),
    ("Gravity",         "F_ij  =  G · m_i · m_j / (|x_i−x_j|² + ε)   for all pairs  i≠j"),
    ("Collision",       "if |x_i−x_j| < r_c:  Δv = elastic_exchange(m_i, m_j, v_i, v_j)"),
    ("Diamond border",  "width(r)  =  w_0 + r·e   [r≤W]     width(r)  =  w_max − (r−W)·c   [r>W]"),
    ("Ball exits",      "if  x_i < left(r)  or  x_i > right(r):  remove ball i"),
    ("Magnet field",    "f(r,x,T)  =  (T−x) · phase_scale(r)      [fan out in W, focus in N]"),
    ("Nail update",     "tilt[r,c] += η · m_i · f(r, x_i, target) · (1 / diameter[r,c])"),
    ("Output",          "score[slot]  =  Σ_i  m_i · 𝟙[x_i ∈ slot_range]"),
    ("Winner",          "ŷ  =  argmax( score )"),
]
for i,(label, formula) in enumerate(math_lines):
    y = 0.85 + i*0.49
    rect(s, 0.35, y, 2.0, 0.42, fill=RGBColor(0x09,0x28,0x42))
    txbox(s, label, 0.45, y+0.06, 1.9, 0.34, size=12, bold=True, color=ACCENT)
    txbox(s, formula, 2.45, y+0.06, 10.4, 0.34, size=12, color=WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — GPU/CPU Mapping
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "10", "GPU / CPU Parallelism Strategy")

for col,(hw, col_fill, ops) in enumerate([
    ("GPU  — Massively Parallel", ACCENT, [
        ("Ball trajectories", "Each ball = one thread. All balls in same row computed in parallel."),
        ("Nail updates (per row)", "All nails in same row are independent. SIMD vectorise across columns."),
        ("Ball-ball gravity", "Sparse distance matrix; compute only within proximity band. O(n·k) not O(n²)."),
        ("Output scoring",    "Parallel reduce: each thread accumulates mass into slot bin."),
        ("Batching",          "Multiple training samples run simultaneously — each sample is one warp."),
    ]),
    ("CPU  — Sequential / SIMD", ACCENT2, [
        ("Vocabulary build",  "Pre-training frequency count. Single pass over corpus. Parallel file reads."),
        ("Nail init",         "Random or frequency-informed tilt init. One-time cost."),
        ("Router diamond",    "Shallower grid. Runs on CPU between GPU specialist calls."),
        ("Checkpointing",     "Nail state saved to disk after each epoch. Nail arrays are just float matrices."),
        ("Validation metrics","Accuracy, ball-mass retention ratio — simple scalar ops."),
    ]),
]):
    x = 0.35 + col*6.6
    rect(s, x, 0.75, 6.3, 0.42, fill=col_fill)
    txbox(s, hw, x+0.1, 0.78, 6.1, 0.36, size=15, bold=True, color=BG)
    for i,(op, desc) in enumerate(ops):
        y = 1.25 + i*1.05
        rect(s, x, y, 6.3, 0.38, fill=RGBColor(0x09,0x28,0x42))
        txbox(s, op, x+0.1, y+0.05, 6.1, 0.3, size=13, bold=True, color=col_fill)
        txbox(s, desc, x+0.1, y+0.42, 6.1, 0.55, size=12, color=LIGHT, wrap=True)

divider(s, 6.75)
txbox(s, "Key advantage:  no activation graph storage (no backprop) → GPU memory is used only for nail state + ball positions",
      0.4, 6.82, 12.4, 0.4, size=13, color=GREEN, italic=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — C# App Architecture
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "11", "C# Application — PRM Engine")

modes = [
    ("🏋  Training", GREEN,
     ["Forward pass per sample", "Nail tilt += magnetic update", "Ball-ball physics active",
      "Periodic checkpoint save"]),
    ("🧪  Test", ACCENT,
     ["Forward pass, no updates", "Accuracy = winner == label", "Ball-mass retention reported",
      "No magnet active"]),
    ("🎛  Tune", ACCENT2,
     ["Like Training, lower η", "Small fine-tune dataset", "Specialist role refinement",
      "Preserve thick-nail routes"]),
    ("✅  Val", LIGHT,
     ["Validation set evaluation", "Per-role retention ratio", "Perplexity equivalent",
      "Outputs confusion matrix"]),
]
for i,(title, col, items) in enumerate(modes):
    x = 0.35 + i*3.25
    rect(s, x, 0.75, 3.1, 0.42, fill=col)
    txbox(s, title, x+0.1, 0.78, 2.9, 0.36, size=14, bold=True, color=BG)
    for j,item in enumerate(items):
        rect(s, x, 1.22+j*0.5, 3.1, 0.44, fill=RGBColor(0x08,0x20,0x32) if j%2==0 else RGBColor(0x05,0x18,0x28))
        txbox(s, "▸  "+item, x+0.1, 1.27+j*0.5, 2.9, 0.36, size=12, color=WHITE)

divider(s, 3.35)
components = [
    "VocabularyBuilder",  "DiamondGrid", "BallSimulator",
    "MagnetField", "PRMEngine", "SpecialistRouter",
]
txbox(s, "Core Components:", 0.4, 3.45, 3.0, 0.4, size=14, bold=True, color=ACCENT)
for i,c in enumerate(components):
    x = 0.35 + (i%3)*4.35
    y = 3.9 + (i//3)*0.72
    rect(s, x, y, 4.1, 0.6, fill=RGBColor(0x09,0x28,0x42), line=ACCENT)
    txbox(s, c, x+0.1, y+0.12, 3.9, 0.38, size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

divider(s, 6.2)
txbox(s, "Inference (all modes):  diamond nails fixed → balls fall freely → winner slot = output.  No magnet, no updates.",
      0.4, 6.28, 12.4, 0.45, size=13, color=GREEN, italic=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Open Questions / Next Steps
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
section_header(s, "12", "Open Questions & Next Steps")

rect(s, 0.3, 0.75, 6.0, 0.38, fill=ACCENT2); txbox(s, "Open Design Questions", 0.4, 0.78, 5.8, 0.32, size=14, bold=True, color=BG)
qs = [
    "Grid depth — does depth = abstraction level?",
    "Nail init — random, uniform, or frequency-informed?",
    "Force decay — how fast does gravity drop off?",
    "Collision elasticity — perfectly elastic or damped?",
    "Training stabilisation — prevent tilt explosion?",
    "Nail diameter — fixed or also learnable?",
    "How many roles cover most language tasks?",
    "Long context — distant tokens' gravity too weak?",
]
for i,q in enumerate(qs):
    txbox(s, f"□  {q}", 0.4, 1.2+i*0.52, 5.8, 0.48, size=13, color=LIGHT)

rect(s, 6.6, 0.75, 6.3, 0.38, fill=GREEN); txbox(s, "Next Steps", 6.7, 0.78, 6.1, 0.32, size=14, bold=True, color=BG)
steps = [
    ("Formalise math model", "Translate all physics to GPU-parallel matrix ops"),
    ("C# PRM prototype", "DiamondGrid + BallSimulator + MagnetField"),
    ("Training mode", "Single forward pass with live nail updates"),
    ("Vocabulary builder", "Frequency-ranked token table from corpus"),
    ("Small-scale test", "Validate next-token prediction on tiny corpus"),
    ("Role specialisation", "Train 2–3 role diamonds, compare output quality"),
    ("GPU acceleration", "ILGPU or compute shaders for ball simulation"),
    ("Benchmark vs Transformer", "Training time & accuracy at equivalent scale"),
]
for i,(title,desc) in enumerate(steps):
    y = 1.2+i*0.52
    txbox(s, f"{'①②③④⑤⑥⑦⑧'[i]}  {title}", 6.7, y, 6.1, 0.26, size=13, bold=True, color=GREEN)
    txbox(s, f"   {desc}", 6.7, y+0.24, 6.1, 0.26, size=12, color=SUBTLE, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Summary
# ─────────────────────────────────────────────────────────────────────────────
s = add_slide(); bg(s)
rect(s, 0, 0, 13.33, 7.5, fill=RGBColor(0x06,0x14,0x22))
txbox(s, "PRM  —  Physical Routing Model", 0.5, 1.0, 12.3, 0.8,
      size=38, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txbox(s, "Tokens are balls.  The model is a grid of nails.  The answer is a magnet.",
      0.5, 1.85, 12.3, 0.55, size=20, color=LIGHT, align=PP_ALIGN.CENTER, italic=True)
divider(s, 2.6)

summary = [
    ("No backpropagation",   "Nails update locally as each ball passes — one forward phase"),
    ("No attention matrix",  "Token relationships emerge from gravitational attraction + collision"),
    ("No positional encoding","Ball entry x-position = sequence position — geometric and free"),
    ("Built-in language prior","Vocabulary slot width = frequency — common words easier to produce"),
    ("Natural sparsity",     "Irrelevant balls fall off diamond edges — no dropout needed"),
    ("Role specialists",     "Diamond geometry encodes cognitive style — generalises across domains"),
]
for i,(title,desc) in enumerate(summary):
    x = 0.5 + (i%2)*6.4
    y = 2.8 + (i//2)*1.35
    rect(s, x, y, 6.1, 1.2, fill=RGBColor(0x0A,0x24,0x3C), line=ACCENT)
    txbox(s, "✦  "+title, x+0.15, y+0.1, 5.8, 0.38, size=14, bold=True, color=ACCENT)
    txbox(s, desc, x+0.15, y+0.5, 5.8, 0.6, size=13, color=LIGHT, wrap=True)

txbox(s, "github.com/sppfizer/thinker", 0.5, 7.0, 12.3, 0.35,
      size=13, color=SUBTLE, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────────────────────────────────────
prs.save(r"C:\src\Claude\Thinker\PRM-Design.pptx")
print("Saved: PRM-Design.pptx")
